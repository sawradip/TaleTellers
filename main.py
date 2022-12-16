import requests
import openai
import wget
import os
# url = "https://whisper.lablab.ai/asr"
# payload={}
# files=[
#   ('audio_file',('orpha_cat.wav',open('orphan_cat.wav','rb'),'audio/mpeg'))
# ]
# response = requests.request("POST", url, data=payload, files=files)
# print(response.json())

OPENAI_KEY = "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"
openai.api_key = OPENAI_KEY

def transcribe(wavPath):
    url = "https://whisper.lablab.ai/asr"
    payload={}
    files=[
    ('audio_file',('orpha_cat.wav',open('orphan_cat.wav','rb'),'audio/mpeg'))
    ]
    response = requests.request("POST", url, data=payload, files=files)
    text = response.json()['text']
    return text

def saveStoryTitleAndParagraphs(story_background, n_paragraphs = 4, model = "text-davinci-003", story_type = "children Story"):
    prompt = f""" Write a {story_type} with {n_paragraphs} paragraphs that takes into account the following backstory:{story_background}
    --
    Formatting: Split paragraphs with extra whitespace between them. Keep the story language child friendly.
    """

    response = openai.Completion.create(
            model=model,
            prompt= prompt,
            max_tokens = 500
        )
    story_text = response["choices"][0]["text"]

    prompt = f"Write a short, catchy, children book appropriate title for the story: {story_text}"

    response = openai.Completion.create(
        model=model,
        prompt= prompt,
        max_tokens = 20
    )
    story_title = response["choices"][0]["text"]

    story_paragraphs = story_text.split('\n')
    story_paragraphs = [para for para in story_paragraphs if len(para)>1]

    return story_title, story_paragraphs



def createDallEPrompt(para_text, model = "text-davinci-003"):
    prompt = f"Create a visually descriptive prompt usable in a children book for a text-to-image AI model based  on the following children's story: {para_text}"
    response = openai.Completion.create(
        model=model,
        prompt= prompt,
        max_tokens = 500
      )
    
    dalle_prompt = response["choices"][0]["text"]
    return dalle_prompt

def saveDallEImage(dallePrompt, n = 0, artistName = "Tasha Tudor", saveDir = "", isTitle = False):
    response = openai.Image.create(
        prompt= "Create a cover image for this story: " if isTitle else "" +  dallePrompt + f"[{artistName} style]",
        n=1,
        size="512x512"
        )
    image_url = response['data'][0]['url']

    outPath = os.path.join(saveDir, f"image_{n}.png")
    wget.download(image_url, out = outPath)
    return outPath

def uniquifyStoryDir(path):
    counter = 1
    pathNew = path
    while os.path.exists(pathNew):
        pathNew = path + "-" + str(counter)
        counter += 1
    return pathNew

uniquifyStoryDir("children Story")

def text2storyfiles(story_background):
    n_paragraph = 4
    model = "text-davinci-003"
    story_type = "children Story"
    # story_style_and_length = f"Children Story with a Title and {n_paragraph} short Paragrapghs"
    # story_background =  f"King Sawradip of Bangladesh commands all of his citizens to be killed. A small boy named Joy is worried for his parents"

    story_title, story_paragraphs = saveStoryTitleAndParagraphs(story_background, 
                                                                n_paragraphs = 4, 
                                                                model = "text-davinci-003", 
                                                                story_type = "children Story")
    
    print("Story Written...")
    storyDir = uniquifyStoryDir(story_type)
    os.mkdir(storyDir)
    print("Story Saved...")
    with open(os.path.join(storyDir, "title.txt"), "w") as f:
        f.write(story_title)

    full_story = ' '.join(story_paragraphs)
    imgPath = saveDallEImage(full_story, n = 0, artistName = "Doraemon", saveDir = storyDir, isTitle = True)

    for(i, para) in enumerate(story_paragraphs):
        with open(os.path.join(storyDir, f"para_{i+1}.txt"), "w") as f:
            f.write(para)

        dalle_prompt = para # createDallEPrompt(para)

        imgPath = saveDallEImage(dalle_prompt, n = i+1, artistName = "Doraemon", saveDir = storyDir)
    print("Images Saved...")

    return storyDir


from glob import glob
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt 
from PIL import Image
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR



# storyDir = "/content/children Story"
def createPPTX(storyDir):
    pre = Presentation()

    title_file = os.path.join(storyDir, "title.txt")
    with open(title_file) as f:
        story_title = f.read()

    # Title Slide
    slideTitle = pre.slides.add_slide(pre.slide_layouts[0])
    p = slideTitle.shapes.title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = story_title
    font = run.font
    font.name = "Comic Sans MS"
    font.size = Pt(25)
    # title = slideTitle.shapes.title
    # title.text_frame.paragraphs[0].font.size = Pt(25)
    # title.text_frame.paragraphs[0].font.name = "Comic Sans MS"
    # title.text = story_title
    # shapes = slideTitle.shapes



    # Remove subtitle
    subtitle = slideTitle.placeholders[1]
    rp = subtitle.element
    rp.getparent().remove(rp)

    # Front page color background
    background = slideTitle.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 255, 255)


    # Story Pages
    para_files = glob(os.path.join(storyDir, "para_*.txt"))
    para_files = sorted(para_files)
    for (i, para_file) in enumerate(para_files):
        with open(para_file) as f:
            para = f.read()
    
        slideNew = pre.slides.add_slide(pre.slide_layouts[8])
        paraTitle = slideNew.shapes.title

        # Background color
        background = slideNew.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(153, 255, 255)


        paraTitle.text = para
        # paraText = paraTitle.text

        # Remove title box
        paraTitle = slideNew.shapes.title
        st = paraTitle.element
        st.getparent().remove(st)

        # Remove subtitle
        paraSubtitle = slideNew.placeholders[2]
        sp = paraSubtitle.element
        sp.getparent().remove(sp)

        # Add new text box for textwrap
        left = Inches(1)
        height = Inches(1.75)
        width = Inches(8)
        top = Inches(5.5)
        txBox = slideNew.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = para

        font = run.font
        font.name = "Comic Sans MS"

        # tf.text = para
        # tf2= tf.shapes.run
        # run = para.add_run()
        # font = run.font
        # tf.font.name = 'Comic-Sans'
        tf.margin_bottom = Inches(0.08)
        tf.margin_left = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        def _add_image(slide, placeholder_id, image_url):
            placeholder = slide.placeholders[placeholder_id]

            # Calculate the image size of the image
            im = Image.open(image_url)
            width, height = im.size

            # Make sure the placeholder doesn't zoom in
            placeholder.height = height
            placeholder.width = width

            # Insert the picture
            placeholder = placeholder.insert_picture(image_url)

            # Calculate the ratio and compare
            image_ratio = width / height
            placeholder_ratio = placeholder.width / placeholder.height
            ratio_difference = placeholder_ratio - image_ratio

            # Placeholder width too wide
            if ratio_difference > 0:
                difference_on_each_side = ratio_difference / 2
                placeholder.crop_left = -difference_on_each_side
                placeholder.crop_right = -difference_on_each_side
            # Placeholder height too high
            else:
                difference_on_each_side = -ratio_difference / 2
                placeholder.crop_bottom = -difference_on_each_side
                placeholder.crop_top = -difference_on_each_side

        image = os.path.join(storyDir, f"image_{i+1}.png")
        _add_image(slideNew, 1, image)


    #end slide
    slideEnd = pre.slides.add_slide(pre.slide_layouts[0])
    run = slideEnd.shapes.title.text_frame.paragraphs[0].add_run()
    run.text = "THE END"
    run.font.name = "Comic Sans MS"

    # story_end = "THE END"
    # slideEnd = pre.slides.add_slide(pre.slide_layouts[0])
    # title = slideEnd.shapes.title
    # subtitle = slideEnd.placeholders[1]
    # title.text = story_end
    # subtitle.text = "See you in another tale."

    # Last page color background
    background = slideEnd.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 255, 255)

    pre.save('StoryBook.pptx')


if __name__ == "__main__":
    story_background = transcribe("orphan_cat.wav")
    print(story_background)
    storyDir = text2storyfiles(story_background)
    createPPTX(storyDir)
